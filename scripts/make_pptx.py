# -*- coding: utf-8 -*-
"""
生成 工研院中分院智慧棒球科技 簡報 PPTX
15 頁，深色主題 + 藍色 (#0047FF) 強調色
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
import copy, os

# ─────────────────────────────────────────────
#  顏色定義
# ─────────────────────────────────────────────
C_DARK      = RGBColor(0x08, 0x08, 0x08)   # 深色背景
C_DARK2     = RGBColor(0x10, 0x10, 0x18)   # 稍淺深色
C_LIGHT     = RGBColor(0xF8, 0xF7, 0xF4)   # 淺色背景
C_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
C_ACCENT    = RGBColor(0x00, 0x47, 0xFF)   # 主題藍
C_ACCENT_L  = RGBColor(0x33, 0x6E, 0xFF)   # 淺藍
C_DARK_TEXT = RGBColor(0x0A, 0x0A, 0x0A)
C_MUTED_L   = RGBColor(0x60, 0x60, 0x60)   # 深色背景上的灰字
C_MUTED_D   = RGBColor(0xA0, 0xA0, 0xA8)   # 淺色背景上的灰字
C_TABLE_H   = RGBColor(0x00, 0x30, 0xCC)   # 表格標題行
C_TABLE_ALT = RGBColor(0xEE, 0xEE, 0xF5)   # 表格交替行

# ─────────────────────────────────────────────
#  字型 (Windows 中文系統常見)
# ─────────────────────────────────────────────
FONT_ZH = "Microsoft JhengHei"   # 微軟正黑體
FONT_EN = "Calibri"

W  = Inches(13.333)   # 16:9 寬
H  = Inches(7.5)      # 16:9 高
M  = Inches(0.55)     # 左右邊距
TM = Inches(0.45)     # 上邊距

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]   # 純空白版型

# ─────────────────────────────────────────────
#  工具函數
# ─────────────────────────────────────────────

def set_bg(slide, color: RGBColor):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def txb(slide, text, l, t, w, h,
        sz=18, bold=False, color=C_WHITE,
        align=PP_ALIGN.LEFT, zh=True, italic=False, space_before=0):
    """在 slide 上加一個單段文字框"""
    shape = slide.shapes.add_textbox(l, t, w, h)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    if space_before:
        p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(sz)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = FONT_ZH if zh else FONT_EN
    # 確保東亞字型也設定
    rPr = run._r.get_or_add_rPr()
    ea = etree.SubElement(rPr, qn('a:ea'))
    ea.set('typeface', FONT_ZH)
    return shape

def bullets(slide, items, l, t, w, h, sz=17, color=C_WHITE,
            indent=True, spacing=10, leading_dot=True):
    """多行 bullet list"""
    shape = slide.shapes.add_textbox(l, t, w, h)
    tf = shape.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(spacing)
        run = p.add_run()
        prefix = "▪  " if leading_dot else ""
        run.text = prefix + item
        run.font.size = Pt(sz)
        run.font.color.rgb = color
        run.font.name = FONT_ZH
        rPr = run._r.get_or_add_rPr()
        ea = etree.SubElement(rPr, qn('a:ea'))
        ea.set('typeface', FONT_ZH)
    return shape

def rect(slide, l, t, w, h, fill=None, line=None, lw=1.5):
    from pptx.util import Pt as _Pt
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    shape = slide.shapes.add_shape(1, l, t, w, h)   # 1 = rectangle
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(lw)
    else:
        shape.line.fill.background()
    return shape

def accent_bar(slide, l, t, w=Inches(0.6), h=Pt(3)):
    return rect(slide, l, t, w, h, fill=C_ACCENT)

def label(slide, text, l, t, w=Inches(6)):
    """章節標籤 — 小灰字"""
    return txb(slide, text, l, t, w, Pt(20),
               sz=10, color=C_MUTED_L, bold=False)

def dark_label(slide, text, l, t, w=Inches(6)):
    return txb(slide, text, l, t, w, Pt(20),
               sz=10, color=C_MUTED_D)

def video_ph(slide, l, t, w, h, filename="", caption=""):
    """影片佔位框"""
    r = rect(slide, l, t, w, h,
             fill=RGBColor(0x10, 0x10, 0x22),
             line=C_ACCENT, lw=1.5)
    # 圖示 + 說明
    cx = l + w // 2
    cy = t + h // 2
    txb(slide, "▶", cx - Inches(0.3), cy - Inches(0.5),
        Inches(0.6), Inches(0.6),
        sz=28, color=C_ACCENT, align=PP_ALIGN.CENTER)
    if filename:
        txb(slide, filename,
            l + Inches(0.2), cy + Inches(0.05),
            w - Inches(0.4), Inches(0.5),
            sz=11, color=C_MUTED_L, align=PP_ALIGN.CENTER)
    if caption:
        txb(slide, caption,
            l + Inches(0.2), cy + Inches(0.5),
            w - Inches(0.4), Inches(0.4),
            sz=10, color=C_MUTED_L, align=PP_ALIGN.CENTER)
    return r

def page_num(slide, text, dark=True):
    c = RGBColor(0x40, 0x40, 0x40) if dark else RGBColor(0xC0, 0xC0, 0xC0)
    txb(slide, text,
        W - Inches(1.5), H - Inches(0.4),
        Inches(1.3), Inches(0.35),
        sz=9, color=c, align=PP_ALIGN.RIGHT)

# ─────────────────────────────────────────────
#  SLIDE 01: 封面
# ─────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
set_bg(s, C_DARK)

# 頂部半白色分隔線
rect(s, M, TM + Inches(3.6), W - M*2, Pt(1), fill=RGBColor(0x28,0x28,0x28))

# 品牌標記
dark_label(s, "工研院中分院 × 智慧棒球科技", M, TM, Inches(8))

# 主標題
txb(s, "場域科技化\n與跨域加值應用", M, TM + Inches(0.5), W - M*2, Inches(2.6),
    sz=52, bold=True, color=C_WHITE)

# Tagline
txb(s, "整合 AI 視覺辨識、無標記式生物力學與雲端數據平台\n打造台灣棒球的科技基礎建設",
    M, TM + Inches(2.8), W - M*2, Inches(0.8),
    sz=14, color=C_MUTED_L)

# 分隔線
accent_bar(s, M, TM + Inches(3.75), Inches(8), Pt(2))

# 三組統計數字（底部列）
col_w = (W - M*2) / 3
for i, (num, lbl) in enumerate([
    ("94%", "球路辨識準確率\n對標 MLB Statcast"),
    ("4+",  "職業場域已導入\n天母·洲際·龍潭·三重"),
    ("14台","高速攝影機\n洲際棒球場部署"),
]):
    cx = M + col_w * i
    txb(s, num, cx, H - Inches(2.4), col_w, Inches(1.0),
        sz=40, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
    txb(s, lbl, cx, H - Inches(1.4), col_w, Inches(0.9),
        sz=13, color=C_MUTED_L, align=PP_ALIGN.CENTER)
    if i < 2:
        rect(s, cx + col_w - Pt(0.5), H - Inches(2.2), Pt(1), Inches(1.8),
             fill=RGBColor(0x28, 0x28, 0x28))

page_num(s, "01 / 15", dark=False)

# ─────────────────────────────────────────────
#  SLIDE 02: 開場影片
# ─────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
set_bg(s, C_DARK)

dark_label(s, "開場影片 · Demo ①", M, TM, Inches(8))
txb(s, "天母棒球場 × 工研院鷹眼系統",
    M, TM + Inches(0.4), W - M*2, Inches(1.0),
    sz=34, bold=True, color=C_WHITE)
txb(s, "全場域運作實景 — 讓數據說話",
    M, TM + Inches(1.3), W - M*2, Inches(0.5),
    sz=15, color=C_MUTED_L)

video_ph(s, M, TM + Inches(1.9), W - M*2, H - TM - Inches(2.5),
         filename="assets/demo1_tianmu_overview.mp4",
         caption="Demo ①  天母棒球場全場域實景")

page_num(s, "02 / 15", dark=False)

# ─────────────────────────────────────────────
#  SLIDE 03: 國際趨勢
# ─────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
set_bg(s, C_LIGHT)

label(s, "第一章 — 為什麼現在需要智慧棒球？", M, TM)
txb(s, "國際趨勢：MLB 已走了 20年",
    M, TM + Inches(0.4), W - M*2, Inches(1.2),
    sz=38, bold=True, color=C_DARK_TEXT)
accent_bar(s, M, TM + Inches(1.55), Inches(0.6), Pt(3))

bullets(s, [
    "美國職棒自 2002 年起導入 AI 數據採集，球員分析、傷病預防已成標配",
    "生物力學應用顯著成果：芝加哥小熊隊導入後，Tommy John 手術比例大幅降低",
    "台灣職棒聯盟規模較小，相關技術導入空白仍大，先發優勢可期",
], M, TM + Inches(1.8), W - M*2, Inches(3.5),
sz=19, color=C_DARK_TEXT, spacing=18)

page_num(s, "03 / 15", dark=True)

# ─────────────────────────────────────────────
#  SLIDE 04: 整體產品架構
# ─────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
set_bg(s, C_LIGHT)

label(s, "第二章 — 我們的核心技術", M, TM)
txb(s, "整體產品架構",
    M, TM + Inches(0.4), W - M*2, Inches(0.9),
    sz=38, bold=True, color=C_DARK_TEXT)
accent_bar(s, M, TM + Inches(1.25), Inches(0.6), Pt(3))

cards = [
    ("01", "電子好球帶",    "Pitch Analysis",
     "AI 攝影機即時球道追蹤，16 項投球指標\n球速 / 旋轉數 / 旋轉軸 / 尾勁全覆蓋"),
    ("02", "無標記式生物力學", "Biomechanics",
     "AI 2D→3D 關節座標重建\n髖肩分離、肘關節內翻力矩，無需穿戴感測器"),
    ("03", "雲端數據平台",  "Cloud Platform",
     "賽後雲端分析、歷史數據查詢\n跨部會資源整合，即時 3D 渲染轉播服務"),
]
card_w = (W - M*2 - Inches(0.3)*2) / 3
for i, (num, title, tag, desc) in enumerate(cards):
    cx = M + (card_w + Inches(0.3)) * i
    cy = TM + Inches(1.55)
    r = rect(s, cx, cy, card_w, H - cy - Inches(0.55),
             fill=C_WHITE, line=RGBColor(0xDD,0xDD,0xE8), lw=1)
    txb(s, num, cx + Inches(0.25), cy + Inches(0.25),
        card_w, Inches(0.7), sz=28, bold=True, color=C_ACCENT)
    txb(s, title, cx + Inches(0.25), cy + Inches(0.9),
        card_w - Inches(0.3), Inches(0.6), sz=18, bold=True, color=C_DARK_TEXT)
    txb(s, desc, cx + Inches(0.25), cy + Inches(1.45),
        card_w - Inches(0.3), Inches(1.8), sz=13, color=C_MUTED_D)
    txb(s, tag, cx + Inches(0.25), cy + Inches(3.3),
        card_w - Inches(0.3), Inches(0.4),
        sz=10, color=C_ACCENT, bold=True)

page_num(s, "04 / 15", dark=True)

# ─────────────────────────────────────────────
#  SLIDE 05: 電子好球帶 (左文 / 右影片)
# ─────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
set_bg(s, C_DARK)

LW = W * 0.46   # 左欄寬
VL = M + LW + Inches(0.4)   # 右欄 left
VW = W - VL - M             # 右欄寬

dark_label(s, "核心技術 ①", M, TM)
txb(s, "電子好球帶",
    M, TM + Inches(0.4), LW, Inches(1.1),
    sz=40, bold=True, color=C_WHITE)
accent_bar(s, M, TM + Inches(1.45), Inches(0.55), Pt(3))

bullets(s, [
    "通用球路辨識，準確率 94%，對標 MLB Statcast",
    "16 項投球指標即時量測：球速、旋轉數、旋轉軸、縫線效應（SSW）",
    "AI 3D 好球帶格線，精準判決輔助",
    "尾勁（Break）計算、球路軌跡全程重建",
], M, TM + Inches(1.65), LW, Inches(4.0),
sz=17, color=C_WHITE, spacing=14)

video_ph(s, VL, TM, VW, H - TM - Inches(0.5),
         filename="assets/demo2_ball_tracking.mp4",
         caption="Demo ②  電子好球帶追蹤實錄")

page_num(s, "05 / 15", dark=False)

# ─────────────────────────────────────────────
#  SLIDE 06: 生物力學分析 (左文 / 右影片)
# ─────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
set_bg(s, C_DARK)

dark_label(s, "核心技術 ②", M, TM)
txb(s, "無標記式\n生物力學分析",
    M, TM + Inches(0.4), LW, Inches(1.6),
    sz=38, bold=True, color=C_WHITE)
accent_bar(s, M, TM + Inches(1.95), Inches(0.55), Pt(3))

bullets(s, [
    "AI 攝影機 2D → 3D 關節座標即時重建",
    "無需穿戴感測器，現場比賽環境即可運作",
    "全身運動學指標：髖肩分離角度、肘關節內翻力矩",
    "出手點辨識（Release Point），精準分析動作變化",
], M, TM + Inches(2.15), LW, Inches(3.5),
sz=17, color=C_WHITE, spacing=14)

video_ph(s, VL, TM, VW, H - TM - Inches(0.5),
         filename="assets/demo3_release_point.mp4",
         caption="Demo ③  出手點辨識實錄")

page_num(s, "06 / 15", dark=False)

# ─────────────────────────────────────────────
#  SLIDE 07: 雲端數據平台 (左文 / 右影片)
# ─────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
set_bg(s, C_DARK)

dark_label(s, "核心技術 ③", M, TM)
txb(s, "雲端數據平台",
    M, TM + Inches(0.4), LW, Inches(1.1),
    sz=40, bold=True, color=C_WHITE)
accent_bar(s, M, TM + Inches(1.45), Inches(0.55), Pt(3))

bullets(s, [
    "賽後雲端分析，球員歷史數據多維度查詢",
    "跨部會資源整合架構，數據互通共享",
    "按時段、對手、場地等條件篩選球員表現",
    "即時 3D 渲染轉播數據加值服務",
], M, TM + Inches(1.65), LW, Inches(4.0),
sz=17, color=C_WHITE, spacing=14)

video_ph(s, VL, TM, VW, H - TM - Inches(0.5),
         filename="assets/demo4_data_query.mp4",
         caption="Demo ④  雲端數據查詢介面")

page_num(s, "07 / 15", dark=False)

# ─────────────────────────────────────────────
#  SLIDE 08: 天母棒球場 (左文 / 右影片)
# ─────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
set_bg(s, C_DARK)

dark_label(s, "第三章 — 實際場域導入成果", M, TM)
txb(s, "天母棒球場\n職業級系統導入",
    M, TM + Inches(0.4), LW, Inches(1.6),
    sz=38, bold=True, color=C_WHITE)
accent_bar(s, M, TM + Inches(1.95), Inches(0.55), Pt(3))

bullets(s, [
    "鷹眼系統 × 工研院研發技術完整整合",
    "電子好球帶、全場球道追蹤同步運作",
    "職業中華職棒比賽實戰應用驗證",
], M, TM + Inches(2.15), LW, Inches(3.0),
sz=17, color=C_WHITE, spacing=16)

video_ph(s, VL, TM, VW, H - TM - Inches(0.5),
         filename="assets/demo1_tianmu_overview.mp4",
         caption="Demo ①  天母棒球場全場域實景")

page_num(s, "08 / 15", dark=False)

# ─────────────────────────────────────────────
#  SLIDE 09: 洲際棒球場 (三支影片橫排)
# ─────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
set_bg(s, C_DARK)

dark_label(s, "場域導入成果 · Demo ⑤⑥⑦", M, TM)
txb(s, "洲際棒球場 — 14 台高速攝影機全場系統",
    M, TM + Inches(0.4), W - M*2, Inches(0.85),
    sz=30, bold=True, color=C_WHITE)
txb(s, "深度學習棒球數據採集 · 3D 好球帶格線 · 多角度輔助判決（每壘 3 台 60FPS Full HD）",
    M, TM + Inches(1.2), W - M*2, Inches(0.45),
    sz=13, color=C_MUTED_L)

# 三欄影片
vids = [
    ("assets/demo5_zhongji_1.mp4", "球道追蹤與好球帶"),
    ("assets/demo6_zhongji_2.mp4", "數據即時顯示"),
    ("assets/demo7_zhongji_3.mp4", "多角度輔助判決"),
]
vw = (W - M*2 - Inches(0.3)*2) / 3
vh = H - TM - Inches(2.0)
for i, (fn, cap) in enumerate(vids):
    vx = M + (vw + Inches(0.3)) * i
    vy = TM + Inches(1.75)
    video_ph(s, vx, vy, vw, vh, filename=fn, caption=cap)

page_num(s, "09 / 15", dark=False)

# ─────────────────────────────────────────────
#  SLIDE 10: 龍潭名人堂 & 三重棒球場
# ─────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
set_bg(s, C_LIGHT)

label(s, "場域導入成果", M, TM)
txb(s, "龍潭名人堂 & 三重棒球場",
    M, TM + Inches(0.4), W - M*2, Inches(0.9),
    sz=36, bold=True, color=C_DARK_TEXT)
accent_bar(s, M, TM + Inches(1.25), Inches(0.6), Pt(3))

half = (W - M*2 - Inches(0.3)) / 2
venues = [
    ("龍潭名人堂", "電子好球帶 + 生物力學同步",
     ["電子好球帶系統正式運作",
      "無標記式生物力學同場整合",
      "投手動作 × 球道數據雙軌分析"]),
    ("三重棒球場", "球道追蹤 Phase 1 & 2",
     ["Phase 1：全場球道追蹤建置完成",
      "Phase 2：打擊分析系統導入中",
      "揮棒軌跡重建、擊球落點分析"]),
]
for i, (name, subtitle, blist) in enumerate(venues):
    cx = M + (half + Inches(0.3)) * i
    cy = TM + Inches(1.55)
    ch = H - cy - Inches(0.55)
    r = rect(s, cx, cy, half, ch,
             fill=C_WHITE, line=C_ACCENT, lw=2)
    rect(s, cx, cy, half, Pt(4), fill=C_ACCENT)   # 頂部藍色條
    txb(s, name, cx + Inches(0.25), cy + Inches(0.2),
        half - Inches(0.3), Inches(0.55),
        sz=20, bold=True, color=C_ACCENT)
    txb(s, subtitle, cx + Inches(0.25), cy + Inches(0.75),
        half - Inches(0.3), Inches(0.5),
        sz=14, bold=True, color=C_DARK_TEXT)
    bullets(s, blist, cx + Inches(0.25), cy + Inches(1.25),
            half - Inches(0.35), ch - Inches(1.4),
            sz=15, color=C_DARK_TEXT, spacing=12)

page_num(s, "10 / 15", dark=True)

# ─────────────────────────────────────────────
#  SLIDE 11: 基層棒球
# ─────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
set_bg(s, C_LIGHT)

label(s, "社會影響力", M, TM)
txb(s, "用科技\n幫助基層棒球",
    M, TM + Inches(0.4), W * 0.45, Inches(1.8),
    sz=40, bold=True, color=C_DARK_TEXT)
accent_bar(s, M, TM + Inches(2.15), Inches(0.6), Pt(3))

txb(s, "科技不應只屬於職業球場。工研院將智慧棒球系統帶入偏遠地區學校，"
       "縮短城鄉資源差距，讓每位年輕球員都能獲得職業級的訓練數據支援。",
    M, TM + Inches(2.35), W - M*2, Inches(0.85),
    sz=14, color=C_MUTED_D)

bullets(s, [
    "紅葉國小：台東偏鄉，棒球發源地之一，首批導入智慧訓練系統",
    "忠孝國中：基層培訓場域，透過數據輔助教練決策",
    "降低訓練門檻，培育台灣下一代棒球人才",
], M, TM + Inches(3.25), W - M*2, Inches(3.0),
sz=17, color=C_DARK_TEXT, spacing=15)

page_num(s, "11 / 15", dark=True)

# ─────────────────────────────────────────────
#  SLIDE 12: 114→115年技術躍進 (表格)
# ─────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
set_bg(s, C_LIGHT)

label(s, "第四章 — 研發歷程與未來規劃", M, TM)
txb(s, "114 年 → 115 年  技術躍進",
    M, TM + Inches(0.4), W - M*2, Inches(0.9),
    sz=34, bold=True, color=C_DARK_TEXT)

# 表格
rows_data = [
    ("比較項目", "114 年度", "115 年度（目標）"),
    ("數據捕捉",
     "投打基礎數據（球速 / 轉速 / 位移）",
     "跑壘、守備、揮棒軌跡、縫線分析、擊球落點"),
    ("核心技術",
     "球體軌跡辨識、投球出手點",
     "縫線效應（SSW）、通用球種辨識、球棒軌跡捕捉"),
    ("應用系統",
     "基礎數據顯示介面",
     "雲地混合架構、3D 即時渲染轉播服務"),
    ("場域推廣",
     "職業球場為主（天母 / 洲際）",
     "名人堂、三重棒球場、基層學校、國際市場"),
]

tl = M
tt = TM + Inches(1.4)
tw = W - M*2
th = H - tt - Inches(0.55)

table = s.shapes.add_table(len(rows_data), 3, tl, tt, tw, th).table
col_widths = [Inches(2.1), Inches(4.5), Inches(5.5)]
for ci, cw in enumerate(col_widths):
    table.columns[ci].width = cw

for ri, row in enumerate(rows_data):
    for ci, cell_text in enumerate(row):
        cell = table.cell(ri, ci)
        cell.text = cell_text
        # 樣式
        p = cell.text_frame.paragraphs[0]
        run = p.runs[0] if p.runs else p.add_run()
        run.text = cell_text
        run.font.name = FONT_ZH
        run.font.size = Pt(14 if ri > 0 else 15)
        run.font.bold = (ri == 0 or ci == 0)
        if ri == 0:
            run.font.color.rgb = C_WHITE
        elif ci == 0:
            run.font.color.rgb = C_DARK_TEXT
        else:
            run.font.color.rgb = C_DARK_TEXT
        # 背景色
        fill = cell.fill
        if ri == 0:
            fill.solid()
            fill.fore_color.rgb = C_TABLE_H
        elif ri % 2 == 0:
            fill.solid()
            fill.fore_color.rgb = C_TABLE_ALT
        else:
            fill.solid()
            fill.fore_color.rgb = C_WHITE
        # 東亞字型
        rPr = run._r.get_or_add_rPr()
        ea = etree.SubElement(rPr, qn('a:ea'))
        ea.set('typeface', FONT_ZH)

page_num(s, "12 / 15", dark=True)

# ─────────────────────────────────────────────
#  SLIDE 13: 完整使用情境 + Demo⑧
# ─────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
set_bg(s, C_DARK)

dark_label(s, "完整使用情境 · Demo ⑧", M, TM)
txb(s, "即時訓練回饋 → 賽後雲端分析 → 教練決策輔助",
    M, TM + Inches(0.4), W - M*2, Inches(0.9),
    sz=28, bold=True, color=C_WHITE)
txb(s, "從牛棚到雲端，完整訓練數據閉環",
    M, TM + Inches(1.25), W - M*2, Inches(0.45),
    sz=15, color=C_MUTED_L)

video_ph(s, M, TM + Inches(1.85), W - M*2, H - TM - Inches(2.4),
         filename="assets/demo8_training_scenario.mp4",
         caption="Demo ⑧  牛棚訓練 → 雲端分析完整流程")

page_num(s, "13 / 15", dark=False)

# ─────────────────────────────────────────────
#  SLIDE 14: 未來方向 + Demo⑨
# ─────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
set_bg(s, C_DARK)

dark_label(s, "第五章 — 未來方向 · Demo ⑨", M, TM)
txb(s, "115 年度\n三大發展重點",
    M, TM + Inches(0.4), LW, Inches(1.6),
    sz=38, bold=True, color=C_WHITE)
accent_bar(s, M, TM + Inches(1.95), Inches(0.55), Pt(3))

bullets(s, [
    "AI 驅動賽訓創新：縫線效應、通用球種辨識模型、3D 轉播渲染",
    "基層 & 海外拓展：名人堂、三重球場全面升級，進軍國際市場",
    "數據加值服務：跑壘守備數據整合，雲地混合架構商業化",
], M, TM + Inches(2.15), LW, Inches(3.5),
sz=17, color=C_WHITE, spacing=18)

video_ph(s, VL, TM, VW, H - TM - Inches(0.5),
         filename="assets/demo9_motion_capture.mp4",
         caption="Demo ⑨  完整動作捕捉實錄")

page_num(s, "14 / 15", dark=False)

# ─────────────────────────────────────────────
#  SLIDE 15: 合作邀請
# ─────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
set_bg(s, C_DARK)

# 大型藍色裝飾塊
rect(s, W - Inches(4.0), 0, Inches(4.0), H,
     fill=RGBColor(0x00, 0x30, 0xCC))

dark_label(s, "合作邀請", M, TM)
txb(s, "一起打造\n台灣棒球的\n科技未來",
    M, TM + Inches(0.5), W * 0.5, Inches(3.2),
    sz=46, bold=True, color=C_WHITE)
txb(s, "工研院中分院具備完整的技術研發能力與場域導入經驗，\n"
       "歡迎各球隊、聯盟、政府單位與我們洽談合作方案。",
    M, TM + Inches(3.7), W * 0.5, Inches(1.0),
    sz=14, color=C_MUTED_L)

# 聯絡資訊
infos = [("單位", "工研院中分院"), ("主題", "智慧棒球科技"), ("年度", "2026")]
for i, (lbl_t, val) in enumerate(infos):
    cy2 = H - Inches(1.5) + Inches(0.42) * i
    txb(s, lbl_t, M, cy2, Inches(1.0), Inches(0.38),
        sz=10, color=C_MUTED_L)
    txb(s, val, M + Inches(1.1), cy2, Inches(3.5), Inches(0.38),
        sz=14, bold=True, color=C_WHITE)

page_num(s, "15 / 15", dark=False)

# ─────────────────────────────────────────────
#  儲存
# ─────────────────────────────────────────────
out = r"C:\Users\User\createslide\智慧棒球科技簡報.pptx"
prs.save(out)
print(f"✅  已儲存：{out}")
print(f"    共 {len(prs.slides)} 頁")
